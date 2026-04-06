import pandas as pd
from pptx import Presentation


def generate_ppt(file_path, output_path="output_report.pptx"):
    try:
        # Read Excel file
        df = pd.read_excel(file_path)

        # Create presentation
        prs = Presentation()

        # -------- Slide 1: Title Slide --------
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        slide.shapes.title.text = "Excel Data Report"
        slide.placeholders[1].text = "Automatically generated from uploaded Excel file"

        # -------- Slide 2: Dataset Overview --------
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = "Dataset Overview"

        content = slide.placeholders[1]

        text = (
            f"Total Rows: {len(df)}\n"
            f"Total Columns: {len(df.columns)}\n\n"
            f"Columns:\n"
            + "\n".join(df.columns)
        )

        content.text = text

        # -------- Slide 3: Sample Data --------
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = "Sample Records"

        sample_text = df.head().to_string(index=False)

        slide.placeholders[1].text = sample_text

        # -------- Slide 4: Statistics --------
        numeric_cols = df.select_dtypes(include="number")

        if not numeric_cols.empty:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            slide.shapes.title.text = "Statistics"

            stats_text = numeric_cols.describe().to_string()

            slide.placeholders[1].text = stats_text

        # Save PPT
        prs.save(output_path)

        return "PowerPoint generated successfully"

    except Exception as e:
        return f"Error generating PPT: {str(e)}"
